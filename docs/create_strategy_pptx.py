#!/usr/bin/env python3
"""Generate Strategy Alignment PowerPoint"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colours
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT_BLUE = RGBColor(0x15, 0x65, 0xC0)
LIGHT_BLUE = RGBColor(0x42, 0xA5, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xCC, 0xCC, 0xCC)
GREEN = RGBColor(0x4C, 0xAF, 0x50)
DARK_TEXT = RGBColor(0x33, 0x33, 0x33)
TABLE_HEADER_BG = RGBColor(0x15, 0x65, 0xC0)
TABLE_ROW_LIGHT = RGBColor(0xF5, 0xF7, 0xFA)
TABLE_ROW_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TABLE_BORDER = RGBColor(0xDD, 0xDD, 0xDD)


def add_dark_background(slide):
    """Add dark background to slide"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG


def add_white_background(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


# ============================================================
# SLIDE 1 — Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_dark_background(slide)

# Title
txBox = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(11.3), Inches(1.2))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Gen3 Firmware Updater App"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# Subtitle
txBox = slide.shapes.add_textbox(Inches(1), Inches(3.0), Inches(11.3), Inches(0.8))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Strategy Alignment"
p.font.size = Pt(28)
p.font.color.rgb = LIGHT_BLUE
p.alignment = PP_ALIGN.CENTER

# Tagline
txBox = slide.shapes.add_textbox(Inches(2), Inches(4.2), Inches(9.3), Inches(1.0))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "How the app delivers against Pure Electric's strategic objectives"
p.font.size = Pt(18)
p.font.color.rgb = LIGHT_GREY
p.alignment = PP_ALIGN.CENTER

# Footer
txBox = slide.shapes.add_textbox(Inches(1), Inches(6.2), Inches(11.3), Inches(0.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "DRAFT  |  February 2026"
p.font.size = Pt(14)
p.font.color.rgb = LIGHT_GREY
p.alignment = PP_ALIGN.CENTER


# ============================================================
# SLIDE 2 — Overview: App's Strategic Role
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_background(slide)

# Title
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.8))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "The App's Strategic Role"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = WHITE

# Quote from strategy
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(11.7), Inches(1.0))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "\"The app reduces risk, increases control, enhances customer experience, and opens up new monetisation channels. It is central to our plan to build a digitally enabled premium brand.\""
p.font.size = Pt(16)
p.font.italic = True
p.font.color.rgb = LIGHT_BLUE
p.alignment = PP_ALIGN.LEFT

# Subtitle
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(2.3), Inches(11.7), Inches(0.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "\u2014 Pure Electric Strategic Plan, Section 7"
p.font.size = Pt(12)
p.font.color.rgb = LIGHT_GREY

# 5 pillars
pillars = [
    ("1. Customer Contact & Safety", "Direct reach to every registered user for recalls and safety notices"),
    ("2. Upsell & Lifecycle Value", "Channel for accessories, servicing, and new scooter offers"),
    ("3. Usage Confirmation", "Confirm customers have read safety instructions before riding"),
    ("4. Leasing & Connected Services", "Scooter lockout for finance defaults, usage-based services"),
    ("5. Technology Integration", "BLE diagnostics, firmware updates, and usage logging"),
]

y_start = 3.1
for i, (title, desc) in enumerate(pillars):
    # Box background
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.8), Inches(y_start + i * 0.85), Inches(11.7), Inches(0.75)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x25, 0x25, 0x40)
    shape.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(1.0), Inches(y_start + i * 0.85 + 0.05), Inches(3.5), Inches(0.35))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = LIGHT_BLUE

    txBox = slide.shapes.add_textbox(Inches(1.0), Inches(y_start + i * 0.85 + 0.38), Inches(11.3), Inches(0.35))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(13)
    p.font.color.rgb = LIGHT_GREY


# ============================================================
# SLIDES 3-7 — Strategy Alignment Details (2 per slide)
# ============================================================
objectives = [
    {
        "objective": "1. Customer Contact & Recall Safety",
        "strategy": "If there is ever an issue with a product, the app gives us direct contact with every registered customer. Registration at setup will be mandatory.",
        "delivered": [
            "Mandatory registration captures email at setup",
            "Push notifications (FCM) reach all users instantly",
            "Notification templates target by role, hardware version, or specific user",
            "Intercom SDK enables direct in-app messaging",
        ]
    },
    {
        "objective": "2. Upsell & Lifecycle Value",
        "strategy": "The app will serve as a channel for offering customers accessories, servicing, and potentially new scooters.",
        "delivered": [
            "In-app Intercom messaging enables direct sales conversations",
            "Push notification system supports promotional campaigns",
            "Audience segmentation by role, hardware version, and region",
            "Foundation for in-app store and accessory catalogue",
        ]
    },
    {
        "objective": "3. Usage Confirmation & Protection",
        "strategy": "The app can confirm that customers have read safety instructions and user manuals before riding. This reduces our liability.",
        "delivered": [
            "Terms & Conditions enforced before app use on every launch",
            "Version-controlled T&C with region and language support",
            "Full audit trail: IP, device info, user agent, time-to-read",
            "Scroll-to-bottom required before acceptance enabled",
        ]
    },
    {
        "objective": "4. Leasing & Connected Services",
        "strategy": "The app will enable features like scooter lockout for leasing or finance contracts where customers default on payments.",
        "delivered": [
            "PIN-secured scooter lock/unlock via BLE",
            "Remote lock capability per scooter",
            "User-scooter ownership tracking in database",
            "Foundation for finance lockout and usage-based services",
        ]
    },
    {
        "objective": "5. Technology Integration",
        "strategy": "The app will connect to scooters via Bluetooth, enabling basic diagnostics, firmware updates, and usage logging.",
        "delivered": [
            "Full BLE telemetry: speed, battery, odometer, range, temperature",
            "Over-the-air firmware updates with progress tracking",
            "Telemetry logged to cloud on every connection",
            "Headlight, cruise control, and lock control via BLE commands",
        ]
    },
    {
        "objective": "6. App-First Support Experience",
        "strategy": "App-first support experience that gives customers troubleshooting tools and issue tracking.",
        "delivered": [
            "Intercom SDK with Fin AI agent for automated support",
            "In-app message composer and conversation history",
            "Self-service support reduces per-incident cost",
            "Scalable across all markets without local support staff",
        ]
    },
    {
        "objective": "7. Reliability & Early Warning System",
        "strategy": "Engineering, customer support, and warranty teams will operate a shared early warning system from field data.",
        "delivered": [
            "Every BLE connection logs firmware, battery health, odometer",
            "Data queryable by hardware version, region, and distributor",
            "Telemetry trends visible in web-admin dashboard",
            "Foundation for automated anomaly detection",
        ]
    },
    {
        "objective": "8. Distributor-Led Global Expansion",
        "strategy": "Distributors are our core global growth channel. We must scale into new markets with minimal capital investment.",
        "delivered": [
            "Web-admin portal for scooter management and firmware uploads",
            "Role-based access: admin, manager, normal user",
            "Region-aware T&C and push notifications",
            "Distributor can manage their fleet independently",
        ]
    },
    {
        "objective": "9. Operational Excellence",
        "strategy": "Use data to drive all decisions. Shift from ticket-handling to system-building: app integration, CRM, real-time diagnostics.",
        "delivered": [
            "Automated firmware update notifications on publish",
            "Notification templates with triggers (firmware, status, scheduled)",
            "Cloud-based scooter registry updated automatically",
            "Edge Functions replace manual processes",
        ]
    },
    {
        "objective": "10. Scale Advantage",
        "strategy": "Our fixed cost base can support significantly more volume with minimal increase in overhead.",
        "delivered": [
            "Near-zero marginal cost: ~0.1p per connection (infrastructure)",
            "Estimated ~£55-95/month platform cost at 100K connections/month",
            "Fin AI support at £0.72/resolution vs £4-11 for human agents",
            "Push notifications (FCM) are free and unlimited at any scale",
        ]
    },
]


def create_detail_slide(obj1, obj2=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_white_background(slide)

    items = [obj1]
    if obj2:
        items.append(obj2)

    for idx, obj in enumerate(items):
        y_offset = 0.3 + idx * 3.5

        # Objective title (blue bar)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(y_offset), Inches(12.3), Inches(0.55)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = ACCENT_BLUE
        shape.line.fill.background()

        txBox = slide.shapes.add_textbox(Inches(0.7), Inches(y_offset + 0.07), Inches(11.9), Inches(0.45))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = obj["objective"]
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = WHITE

        # Two columns
        # Left: Strategy says
        left_x = 0.5
        right_x = 6.8

        txBox = slide.shapes.add_textbox(Inches(left_x), Inches(y_offset + 0.7), Inches(5.8), Inches(0.35))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = "STRATEGY SAYS"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = ACCENT_BLUE

        txBox = slide.shapes.add_textbox(Inches(left_x), Inches(y_offset + 1.0), Inches(5.8), Inches(2.2))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = obj["strategy"]
        p.font.size = Pt(13)
        p.font.italic = True
        p.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

        # Right: App delivers
        txBox = slide.shapes.add_textbox(Inches(right_x), Inches(y_offset + 0.7), Inches(6.0), Inches(0.35))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = "APP DELIVERS"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = GREEN

        for j, item in enumerate(obj["delivered"]):
            txBox = slide.shapes.add_textbox(Inches(right_x), Inches(y_offset + 1.0 + j * 0.45), Inches(6.0), Inches(0.4))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = "\u2713  " + item
            p.font.size = Pt(13)
            p.font.color.rgb = DARK_TEXT


# Create slides 2 objectives at a time
for i in range(0, len(objectives), 2):
    obj2 = objectives[i + 1] if i + 1 < len(objectives) else None
    create_detail_slide(objectives[i], obj2)


# ============================================================
# SLIDE 8 — Summary Table
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_white_background(slide)

txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.6))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Summary: Strategy Alignment at a Glance"
p.font.size = Pt(26)
p.font.bold = True
p.font.color.rgb = ACCENT_BLUE

# Table
rows = 11
cols = 3
tbl = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(0.9), Inches(12.3), Inches(6.2)).table

tbl.columns[0].width = Inches(0.5)
tbl.columns[1].width = Inches(4.3)
tbl.columns[2].width = Inches(7.5)

# Header row
headers = ["#", "Strategy Objective", "How the App Delivers"]
for i, h in enumerate(headers):
    cell = tbl.cell(0, i)
    cell.text = h
    cell.fill.solid()
    cell.fill.fore_color.rgb = TABLE_HEADER_BG
    for para in cell.text_frame.paragraphs:
        para.font.size = Pt(12)
        para.font.bold = True
        para.font.color.rgb = WHITE

summary_data = [
    ("1", "Customer Contact & Recall Safety", "Mandatory registration, FCM push notifications, Intercom messaging"),
    ("2", "Upsell & Lifecycle Value", "In-app messaging, segmented push campaigns, promotion infrastructure"),
    ("3", "Usage Confirmation & Protection", "Enforced T&C acceptance, audit trail, region/language support"),
    ("4", "Leasing & Connected Services", "PIN-secured lock/unlock, ownership tracking, lockout foundation"),
    ("5", "Technology Integration", "BLE telemetry, OTA firmware updates, scooter controls"),
    ("6", "App-First Support", "Intercom + Fin AI, in-app messaging, scalable self-service"),
    ("7", "Reliability Early Warning", "Cloud telemetry logging, queryable by HW version/region"),
    ("8", "Distributor Global Expansion", "Web-admin portal, role-based access, region-aware features"),
    ("9", "Operational Excellence", "Automated notifications, cloud scooter registry, Edge Functions"),
    ("10", "Scale Advantage", "~£55-95/mo at 100K connections; near-zero marginal cost; Fin AI £0.72/resolution"),
]

for r, (num, obj, delivery) in enumerate(summary_data):
    row_idx = r + 1
    bg = TABLE_ROW_LIGHT if r % 2 == 0 else TABLE_ROW_WHITE

    for c, val in enumerate([num, obj, delivery]):
        cell = tbl.cell(row_idx, c)
        cell.text = val
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
        for para in cell.text_frame.paragraphs:
            para.font.size = Pt(11)
            para.font.color.rgb = DARK_TEXT
            if c == 0:
                para.alignment = PP_ALIGN.CENTER
                para.font.bold = True


# ============================================================
# SLIDE 9 — Status & Next Steps
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_background(slide)

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.8))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Current Status & Roadmap Alignment"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = WHITE

# Built & Live
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.4))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "BUILT & LIVE"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = GREEN

live_items = [
    "BLE connection, telemetry & scooter controls",
    "Over-the-air firmware updates",
    "User registration & email verification",
    "Terms & Conditions enforcement with audit trail",
    "PIN-secured scooter lock/unlock",
    "Push notifications with segmented targeting",
    "Intercom support with Fin AI agent",
    "Web-admin portal for distributors",
    "Secure Edge Function API (no direct DB writes)",
]

for i, item in enumerate(live_items):
    txBox = slide.shapes.add_textbox(Inches(1.0), Inches(2.0 + i * 0.45), Inches(5.5), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "\u2713  " + item
    p.font.size = Pt(13)
    p.font.color.rgb = LIGHT_GREY

# Roadmap alignment
txBox = slide.shapes.add_textbox(Inches(7.0), Inches(1.5), Inches(5.5), Inches(0.4))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "ROADMAP ALIGNMENT"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = LIGHT_BLUE

roadmap = [
    ("Next 12 months", [
        "First public release",
        "Distributor & support flow integration",
        "Accessory upsell channel",
    ]),
    ("12\u201324 months", [
        "In-app upselling",
        ">70% app registration rate",
        "Leasing model foundation",
    ]),
    ("24\u201336 months", [
        "Embedded services rollout",
        "Brand leadership in 3+ markets",
        "Usage-based insurance partnerships",
    ]),
]

y = 2.0
for phase, items in roadmap:
    txBox = slide.shapes.add_textbox(Inches(7.2), Inches(y), Inches(5.3), Inches(0.35))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = phase
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = LIGHT_BLUE
    y += 0.35

    for item in items:
        txBox = slide.shapes.add_textbox(Inches(7.4), Inches(y), Inches(5.1), Inches(0.35))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = "\u2022  " + item
        p.font.size = Pt(12)
        p.font.color.rgb = LIGHT_GREY
        y += 0.33
    y += 0.2


# ============================================================
# SLIDE 10 — Platform Cost Breakdown
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_white_background(slide)

txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.6))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Platform Costs at Scale: 100,000 Connections/Month (GBP)"
p.font.size = Pt(26)
p.font.bold = True
p.font.color.rgb = ACCENT_BLUE

# Subtitle
txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(12.3), Inches(0.4))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Usage-based pricing with near-zero marginal cost per connection. Converted at $1 = £0.73 (Feb 2026). Services billed in USD."
p.font.size = Pt(14)
p.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

# Cost table
cost_rows = 9
cost_cols = 4
cost_tbl = slide.shapes.add_table(cost_rows, cost_cols, Inches(0.5), Inches(1.4), Inches(12.3), Inches(4.2)).table

cost_tbl.columns[0].width = Inches(2.8)
cost_tbl.columns[1].width = Inches(4.0)
cost_tbl.columns[2].width = Inches(3.5)
cost_tbl.columns[3].width = Inches(2.0)

cost_headers = ["Service", "What It Covers", "Usage at 100K Connections/Mo", "Est. Monthly Cost"]
for i, h in enumerate(cost_headers):
    cell = cost_tbl.cell(0, i)
    cell.text = h
    cell.fill.solid()
    cell.fill.fore_color.rgb = TABLE_HEADER_BG
    for para in cell.text_frame.paragraphs:
        para.font.size = Pt(11)
        para.font.bold = True
        para.font.color.rgb = WHITE

cost_data = [
    ("Supabase Pro", "Database, Edge Functions (2M incl.), 8GB storage", "~400K Edge Function calls, ~100K telemetry rows", "£18"),
    ("Supabase Compute", "Dedicated Postgres instance", "May need upgrade at scale", "£0 - £37"),
    ("Supabase Storage Overage", "DB storage beyond 8GB at £0.09/GB", "~1.5GB/yr growth; 8GB lasts 5+ years", "£0 - £7"),
    ("SendGrid Essentials", "Verification & notification emails", "~10K emails (not every connection = new user)", "£0 - £15"),
    ("Firebase FCM", "Push notifications to all users", "Unlimited - free at any scale", "£0"),
    ("Intercom Essential", "In-app support + Fin AI agent", "1 seat + £0.72 per AI resolution", "£21 + usage"),
    ("", "", "ESTIMATED TOTAL (excl. Fin AI resolutions)", "£55 - £95/mo"),
]

for r, (svc, covers, usage, cost) in enumerate(cost_data):
    row_idx = r + 1
    bg = TABLE_ROW_LIGHT if r % 2 == 0 else TABLE_ROW_WHITE
    if r >= 5:
        bg = TABLE_ROW_WHITE

    is_total_row = (r == len(cost_data) - 1)
    for c, val in enumerate([svc, covers, usage, cost]):
        cell = cost_tbl.cell(row_idx, c)
        cell.text = val
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
        for para in cell.text_frame.paragraphs:
            para.font.size = Pt(11)
            para.font.color.rgb = DARK_TEXT
            if is_total_row:
                para.font.bold = True
                if c == 3:
                    para.font.color.rgb = ACCENT_BLUE
                    para.font.size = Pt(13)

# Key insights below table
insights = [
    ("Cost per connection:", "~0.1p (infrastructure only) or ~1.1p including AI support"),
    ("Fin AI support cost:", "At 2% contact rate with 70% AI resolution = ~1,400 resolutions/mo = ~£1,020"),
    ("Compared to human agents:", "Fin AI at £0.72/resolution vs £4-11 per human-handled ticket across multiple languages"),
    ("Database storage:", "~750 bytes/telemetry row; 8GB included covers 5+ years at 100K connections/mo"),
]

y_insight = 5.8
for label, detail in insights:
    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(y_insight), Inches(11.9), Inches(0.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = label + "  "
    run.font.size = Pt(12)
    run.font.bold = True
    run.font.color.rgb = DARK_TEXT
    run2 = p.add_run()
    run2.text = detail
    run2.font.size = Pt(12)
    run2.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
    y_insight += 0.42


# Save
output_path = "/Users/catherineives/AndroidStudioProjects/Gen3FirmwareUpdater/docs/App_Strategy_Alignment.pptx"
prs.save(output_path)
print(f"Saved to {output_path}")
